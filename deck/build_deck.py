"""
Builds the Part 1 strategy deck (max 4 slides) as an actual .pptx file.
Run once to generate; content lives here, not hand-edited in PowerPoint,
so it stays in sync with the numbers verified in this same repo.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT_PATH = Path(__file__).parent / "Example Co. Strategy Deck.pptx"
DIAGRAM_PNG = Path(__file__).parent / "architecture_diagram.png"

NAVY = RGBColor(0x1a, 0x1a, 0x2e)
DARK = RGBColor(0x1a, 0x1a, 0x1a)
MUTED = RGBColor(0x6b, 0x72, 0x80)
BLUE = RGBColor(0x1d, 0x4e, 0xd8)
GREEN = RGBColor(0x0f, 0x76, 0x6e)
AMBER = RGBColor(0xb4, 0x53, 0x09)
LIGHT_BG = RGBColor(0xf7, 0xf7, 0xf8)
BORDER = RGBColor(0xd4, 0xd4, 0xd8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_title(slide, kicker, title):
    k = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.4))
    tf = k.text_frame
    tf.text = kicker
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = MUTED
    tf.paragraphs[0].font.bold = True

    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.7))
    tf = t.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.bold = True
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.35), Inches(12.3), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    return slide


def bullet_box(slide, x, y, w, h, heading, items, heading_color=NAVY, font_size=13):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = heading_color
    for item in items:
        p = tf.add_paragraph()
        p.text = "  •  " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)
    return box


def caveat_banner(slide, y, text, height=0.55):
    box = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.3), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xfe, 0xf6, 0xe7)
    box.line.color.rgb = RGBColor(0xfb, 0xbf, 0x24)
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = AMBER
    return box


def stat_card(slide, x, y, w, h, value, label, color=NAVY):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# Slide 1: Operating model and 12-month roadmap
# ---------------------------------------------------------------------------
s1 = add_slide()
add_title(s1, "1 / 4 · OPERATING MODEL", "One routing model, one AI layer, phased rollout")

bullet_box(s1, 0.5, 1.6, 6.0, 2.6, "The unifying idea", [
    "Service, Success, and Sales already share limited intake structure – Support and Sales have real forms, but Success traffic likely lands in a plain, unstructured mailbox with no form fields at all",
    "Route on need, opportunity, and risk detected in the message content – not on which form it arrived through",
    "A single primary owner is always kept per message; other teams with an independent signal are looped in, never left blind",
    "Escalate to a person whenever the system is actually uncertain, rather than forcing a guess",
])

bullet_box(s1, 6.8, 1.6, 6.0, 2.6, "Build / change / remove", [
    "BUILD: the AI triage + loop-in layer (this submission, working end to end)",
    "CHANGE: routing from channel-based to need/opportunity/risk-based",
    "REMOVE: nothing customer-facing in the early phases – the highest-risk step (no human in the loop) is sequenced last, not first",
])

road_y = 4.3
roadmap = [
    ("Phase 1", "Build + test at scale", "Where we are now: synthetic prototype, being re-validated against real tickets at increasing scale (10 → 50 → 100 → 1,000) before any production exposure."),
    ("Phase 2", "Monitored production pilot", "Small, closely-watched slice of real inbound traffic. Every draft still human-approved. Go/no-go on accuracy before expanding."),
    ("Phase 3", "Broader rollout", "Expand across full Service volume once the pilot confirms accuracy holds on real data; extend loop-in/expansion-signal capture to Success."),
    ("Phase 4", "Deepen integration", "Connect to the company's own live systems for real data in replies (tracking/rates/duty/labels); lightweight weekly CSM digest."),
    ("Phase 5", "Evaluate self-service deflection", "Only once production accuracy has a long track record. The one step that removes a human from the loop entirely – sequenced last on purpose, evaluated, not committed to."),
]
col_w = 2.34
for i, (q, head, body) in enumerate(roadmap):
    x = 0.5 + i * (col_w + 0.12)
    card = s1.shapes.add_shape(1, Inches(x), Inches(road_y), Inches(col_w), Inches(2.1))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xea, 0xf2, 0xfd) if i == 0 else LIGHT_BG
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.75)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = q
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.text = head
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.space_after = Pt(4)
    p3 = tf.add_paragraph()
    p3.text = body
    p3.font.size = Pt(9.5)
    p3.font.color.rgb = DARK

caveat_banner(s1, 6.55,
    "Status: early prototype, validated on 120 synthetic messages only. Before any production claim: re-validate against real tickets at increasing scale (10 → 50 → 100 → 1,000) and track accuracy at each step.")

s1.notes_slide.notes_text_frame.text = (
    "Integration detail (for Q&A): the triage agent runs as a background worker subscribed to helpdesk/CRM webhooks "
    "(ticket created / updated) - not a polling loop. It writes its output (category, confidence, draft, flags) back onto "
    "the ticket as a private note/custom field, visible inline to whoever picks it up, using whatever app or macro surface "
    "the helpdesk exposes for that - no separate tool, no context switch. Because it's async, a slow or failed AI call never "
    "blocks a human from working the ticket manually; it only means that ticket's enrichment arrives late or not at all, "
    "which the pipeline already treats as a safe fallback (see error handling). API calls use explicit timeouts and "
    "retry-with-backoff (3 retries, 60s timeout) rather than SDK defaults. An eval-as-CI suite (run_eval.py, wired as a "
    "GitHub Action) runs a fixed set of known-answer messages on every push, so a prompt change that regresses a "
    "previously-fixed bug fails the build automatically rather than being caught by chance."
)

# ---------------------------------------------------------------------------
# Slide 2: AI-first tooling landscape
# ---------------------------------------------------------------------------
s2 = add_slide()
add_title(s2, "2 / 4 · AI-FIRST TOOLING LANDSCAPE", "One stack, purpose-built rows completed against the brief's template")

rows = [
    ("Customer Service (Reactive)", "Helpdesk platform (existing)", "Why: assume this is already the company's system of record for support - the AI layer sits on top of it, avoiding a costly migration the budget doesn't support", "CRM, Analytics, AI triage layer", "CSAT, FRT, Resolution time"),
    ("Customer Success (Proactive)", "CRM/CSM platform (existing tool extended; dedicated tool tbc)", "Why: at a representative ARR/CSM book, a dedicated CSM-platform licence is hard to justify against typical non-staff budget headroom - extending an already-licensed CRM is the lower-cost path unless real usage shows it can't scale", "Finance, Helpdesk, Analytics", "NRR, Churn %, Expansion revenue"),
    ("Inbound Sales", "Sales CRM (existing)", "Why: assume this is already licensed and already driving inbound lead volume - the AI layer routes into this system rather than duplicating it", "Marketing, CS, Helpdesk", "Conversion %, Revenue per lead"),
    ("Omnichannel Layer", "Helpdesk channels + AI routing (this build)", "Why: the existing helpdesk already covers the channels in use - a separate omnichannel platform would be new spend for a gap that doesn't exist; the AI routing layer is the actual gap being filled", "Helpdesk + CRM", "% routed by AI, Wait time"),
    ("Automation & AI", "Claude triage agent (built, prototype)", "Why this architecture, not an off-the-shelf bot: company-specific guardrails (sensitive-topic/retention overrides, multi-team loop-in) need to be auditable and tunable, which a generic vendor tool wouldn't give", "Helpdesk, CSM, CRM", "Cost/inquiry: $0.0103 measured (prototype)"),
    ("Knowledge Management", "Helpdesk KB + AI search (future)", "Why: assume the company already publishes real help-centre content - self-serve AI search should point at what exists, not a rebuilt knowledge base", "Helpdesk, CRM", "KB usage %, Ticket reduction"),
    ("Reporting & Analytics", "BI/dashboarding (this build's dashboard as an early example)", "Why: Service, Success, and Sales don't currently share one reporting view - even a lightweight shared dashboard is a real change from three siloed ones", "All systems", "Revenue impact, Cost-to-serve, NPS"),
]
headers = ["Function", "Core tooling", "Why included", "Integrations needed", "Metrics to track"]
col_widths = [1.6, 2.3, 3.9, 1.7, 2.8]

table_top = Inches(1.55)
table_height = Inches(5.5)
table_shape = s2.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), table_top, Inches(12.3), table_height)
table = table_shape.table
for i, w in enumerate(col_widths):
    table.columns[i].width = Inches(w)

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

for i, row in enumerate(rows, start=1):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff) if i % 2 else LIGHT_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9.5)
            p.font.color.rgb = DARK
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
        cell.vertical_anchor = 3  # middle

footnote2 = s2.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4))
tf = footnote2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Illustrative tooling assumptions for this demo - in a real engagement, each row would be confirmed against the company's actual stack rather than assumed. Dedicated CSM platform (e.g. Gainsight) is a common open question at this stage."
p.font.size = Pt(9)
p.font.color.rgb = MUTED

# ---------------------------------------------------------------------------
# Slide 3: Commercial impact
# ---------------------------------------------------------------------------
s3 = add_slide()
add_title(s3, "3 / 4 · COMMERCIAL IMPACT", "Illustrative estimates from stated assumptions - directionally useful, not final")

# Metric 1
bullet_box(s3, 0.5, 1.55, 4.0, 0.3, "Service cost per inquiry", [], heading_color=BLUE)
stat_card(s3, 0.5, 1.95, 1.9, 1.0, "$5.83", "baseline / inquiry", MUTED)
stat_card(s3, 2.55, 1.95, 1.9, 1.0, "$4.96", "with AI (-15%)", GREEN)
box1 = s3.shapes.add_textbox(Inches(0.5), Inches(3.05), Inches(4.0), Inches(1.85))
tf = box1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "28 FTEs × $15K = $420K/yr ÷ 72,000 inquiries/yr."
p.font.size = Pt(10)
p.font.color.rgb = DARK
p2 = tf.add_paragraph()
p2.text = "Deliberately conservative assumption: AI triage+draft cuts average handling time by 15% (a modest starting estimate, not a target - real figure needs piloting)."
p2.font.size = Pt(10)
p2.font.color.rgb = MUTED
p2.space_before = Pt(6)
p3 = tf.add_paragraph()
p3.text = "Same 28 FTEs then absorb ~18% more volume; AI cost at that volume (real measured rate) stays trivial either way."
p3.font.size = Pt(10)
p3.font.color.rgb = MUTED
p3.space_before = Pt(4)

# Metric 2
bullet_box(s3, 4.85, 1.55, 4.0, 0.3, "Expansion revenue per CSM", [], heading_color=GREEN)
stat_card(s3, 4.85, 1.95, 1.9, 1.0, "$30,000", "baseline / CSM / yr", MUTED)
stat_card(s3, 6.9, 1.95, 1.9, 1.0, "+$6,480", "new (+22%)", GREEN)
box2 = s3.shapes.add_textbox(Inches(4.85), Inches(3.05), Inches(4.0), Inches(1.85))
tf = box2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "5 CSMs × $250K ARR/qtr book; NRR 103% = $30K net expansion/CSM/yr today."
p.font.size = Pt(10)
p.font.color.rgb = DARK
p2 = tf.add_paragraph()
p2.text = "Measured in this build's test set: 6% of messages carry an expansion signal Success wouldn't otherwise own. Heavily discounted to 3% for real traffic (test set is edge-case-enriched, not representative)."
p2.font.size = Pt(10)
p2.font.color.rgb = MUTED
p2.space_before = Pt(6)
p3 = tf.add_paragraph()
p3.text = "3% of 72,000 inquiries ÷ 5 CSMs = 432 signals/CSM/yr. Conservatively assume 5% are actionable, $300 average incremental ARR each → +$6,480/CSM/yr."
p3.font.size = Pt(10)
p3.font.color.rgb = MUTED
p3.space_before = Pt(4)

# Metric 3
bullet_box(s3, 9.2, 1.55, 3.6, 0.3, "Inbound lead conversion efficiency", [], heading_color=AMBER)
stat_card(s3, 9.2, 1.95, 1.7, 1.0, "$400K", "ARR / sales FTE / yr", MUTED)
stat_card(s3, 11.05, 1.95, 1.7, 1.0, "$412K", "with AI (+3%)", AMBER)
box3 = s3.shapes.add_textbox(Inches(9.2), Inches(3.05), Inches(3.6), Inches(1.85))
tf = box3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "3 FTEs closing ~$100K ARR/qtr each = $400K/yr/FTE today."
p.font.size = Pt(10)
p.font.color.rgb = DARK
p2 = tf.add_paragraph()
p2.text = "Modest assumption: correct routing of Sales-flavored content plus a faster AI-drafted first touch lifts conversion on the same lead volume by ~3% - a conservative starting point, not measured live."
p2.font.size = Pt(10)
p2.font.color.rgb = MUTED
p2.space_before = Pt(6)
p3 = tf.add_paragraph()
p3.text = "= +$12K/FTE/yr, +$36K/yr across the 3-person team – same headcount, same leads, a little more converted."
p3.font.size = Pt(10)
p3.font.color.rgb = MUTED
p3.space_before = Pt(4)

budget_note = s3.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.5))
tf = budget_note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Budget reality check"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = "Staff costs already consume $585K of the $600K budget (30×$15K Service + 3×$15K Sales + 5×$18K CS), leaving $15K/yr non-staff headroom. Real measured AI cost at combined-run rates scales to roughly $56–$60/month – well inside that headroom. The framing throughout: AI cost is netted against reallocated human time it frees up, not counted as an extra cost in isolation – humans were doing 100% of this work before; the question is how much of that time AI can take on more cheaply."
p2.font.size = Pt(11)
p2.font.color.rgb = DARK
p2.space_before = Pt(4)

caveat_banner(s3, 6.75,
    "These three figures are calculation chains built on named, stated assumptions (handling-time reduction, signal-conversion rates, response-time lift) - illustrative hypotheses to test, not validated forecasts. No current baseline was measured for any of the three metrics; real deltas require piloting against actual historical data before being relied upon.")

# ---------------------------------------------------------------------------
# Slide 4: Team, skills, and adoption
# ---------------------------------------------------------------------------
s4 = add_slide()
add_title(s4, "4 / 4 · TEAM, SKILLS, ADOPTION", "The skill shift is editing judgement, not engineering")

bullet_box(s4, 0.5, 1.6, 6.0, 3.2, "Skills the team will need", [
    "A dedicated Ops Lead role – not an engineer, but someone who owns testing, monitoring, and running this day to day (config/guardrail tuning, watching for false positives, the eval suite); Head of Client Services provides strategic oversight, not hands-on operation",
    "Guardrail/config literacy, not ML engineering – the whole routing model lives in one config file; owning it needs judgement about thresholds and edge cases, not a data science background",
    "Editing and reviewing judgement over composition – the job shifts from “write a reply” to “approve, edit, or escalate a draft,” a distinctly different skill worth naming and training for explicitly",
    "Cross-team trust in the loop-in mechanism – Sales has to trust Success ownership of a looped-in expansion signal, and vice versa, rather than each team wanting to own everything that touches it",
], font_size=11.5)

bullet_box(s4, 6.8, 1.6, 6.0, 3.2, "How adoption actually happens", [
    "Lead with the “knows what it doesn't know” behavior – reps seeing the system defer to Team Lead Triage instead of confidently guessing wrong builds trust faster than any accuracy statistic quoted at them",
    "Surface guardrail flags and loop-ins inline in the helpdesk, where reps already work – not a separate dashboard they have to remember to check",
    "Run a bake-off first: AI drafts shown but optional, tracked against manual-only tickets, before review-and-send becomes the standard workflow",
    "Sequence trust by risk: Service first (highest volume, lowest political sensitivity), Success/Sales ownership questions only once Service is proven",
    "SLAs apply to both queues, same as today: IC response at volume, and Team Lead/Principal response on escalated cases – this build doesn't introduce a new SLA obligation, it routes into the ones that already exist",
], font_size=11.5)

note = s4.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.5))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Why this order, not the reverse"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = "Judgement, here, is choosing what NOT to do in the early phases: no new seat-based platform purchase (the budget doesn't support it), no customer-facing self-service bot until triage accuracy has a long track record on real production data, and no headcount changes forced by day-one automation. The AI layer is framed throughout as making the existing team faster and more visible to each other, not as a replacement for the judgement already in the room. Team Lead Triage volume (currently ~11-13% of messages) is a maturity signal to track down over time as accuracy improves, not a fixed cost - and future iterations should test Opus, Haiku, and Sonnet against larger real-ticket sets rather than assuming today's model choice holds at scale."
p2.font.size = Pt(10.5)
p2.font.color.rgb = DARK
p2.space_before = Pt(4)

prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH}")
